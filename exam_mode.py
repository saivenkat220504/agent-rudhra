import streamlit as st
import tempfile
import os
import time
import pandas as pd
from PyPDF2 import PdfReader
from openai import OpenAI
import random
import docx
from pptx import Presentation

from score import evaluate_image_answer


# ---------- CONFIG ----------
LLM_API_KEY = os.getenv("OPENAI_API_KEY")
MODEL_NAME = "openai/gpt-4o-mini"

client = OpenAI(
    api_key=LLM_API_KEY,
    base_url="https://openrouter.ai/api/v1"
)


# ==========================================================
# RESET EXAM
# ==========================================================

def _reset_exam():

    keys = [
        "questions",
        "mcq_questions",
        "current_q_idx",
        "results",
        "submitted",
        "start_time",
        "eval_mode",
        "time_limit",
        "uploaded_files",
        "total_pages"
    ]

    for k in keys:
        if k in st.session_state:
            del st.session_state[k]

    st.session_state.exam_step = "upload"


# ==========================================================
# HELPERS
# ==========================================================

def count_pages_in_files(uploaded_files):

    total_pages = 0

    for file in uploaded_files:

        file_type = file.name.split('.')[-1].lower()
        file.seek(0)

        if file_type == "pdf":

            try:
                reader = PdfReader(file)
                total_pages += len(reader.pages)
            except:
                pass

        elif file_type in ["docx","doc"]:

            document = docx.Document(file)

            full_text = "\n".join(
                [p.text for p in document.paragraphs if p.text.strip()]
            )

            pages = max(1,(len(full_text)+2999)//3000)
            total_pages += pages

        elif file_type in ["pptx","ppt"]:

            prs = Presentation(file)
            total_pages += len(prs.slides)

    return total_pages


# ==========================================================
# QUESTION GENERATION
# ==========================================================

def _extract_all_pages(uploaded_files):
    """Common helper to extract text pages from uploaded files."""
    all_pages = []

    for file in uploaded_files:

        file_type = file.name.split('.')[-1].lower()
        file.seek(0)

        if file_type == "pdf":

            try:
                reader = PdfReader(file)

                for page in reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        all_pages.append(text.strip())

            except:
                pass

        elif file_type in ["docx","doc"]:

            document = docx.Document(file)

            full_text = "\n".join(
                [p.text for p in document.paragraphs if p.text.strip()]
            )

            chunks = [
                full_text[i:i+3000]
                for i in range(0,len(full_text),3000)
            ]

            all_pages.extend(chunks)

        elif file_type in ["pptx","ppt"]:

            prs = Presentation(file)

            for slide in prs.slides:

                slide_text=""

                for shape in slide.shapes:

                    if hasattr(shape,"text") and shape.text:
                        slide_text+=shape.text+"\n"

                if slide_text.strip():
                    all_pages.append(slide_text.strip())

    return all_pages


def extract_questions_from_files(uploaded_files, num_questions):

    all_pages = _extract_all_pages(uploaded_files)

    if not all_pages:
        return []

    selected_pages = random.sample(
        all_pages,
        min(num_questions,len(all_pages))
    )

    questions=[]

    for idx,page_text in enumerate(selected_pages,start=1):

        prompt=f"""
Generate ONE clear exam question from the following content.

Content:
{page_text[:3000]}

Format:
{idx}. Question
"""

        response=client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role":"user","content":prompt}],
            temperature=0.2,
            max_tokens=80
        )

        question=response.choices[0].message.content.strip()

        if not question.startswith(str(idx)):
            question=f"{idx}. {question}"

        questions.append(question)

    return questions


# ==========================================================
# MCQ QUESTION GENERATION
# ==========================================================

def extract_mcq_questions_from_files(uploaded_files, num_questions):
    """Generate MCQ questions with 4 options and a correct answer."""

    all_pages = _extract_all_pages(uploaded_files)

    if not all_pages:
        return []

    selected_pages = random.sample(
        all_pages,
        min(num_questions, len(all_pages))
    )

    mcq_list = []

    for idx, page_text in enumerate(selected_pages, start=1):

        prompt = f"""
Generate ONE multiple-choice question (MCQ) from the following content.

Content:
{page_text[:3000]}

You MUST return EXACTLY this format (no extra text):
Q: <question text>
A) <option A>
B) <option B>
C) <option C>
D) <option D>
CORRECT: <letter>
"""

        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=250
        )

        raw = response.choices[0].message.content.strip()
        mcq = _parse_mcq(raw, idx)

        if mcq:
            mcq_list.append(mcq)

    return mcq_list


def _parse_mcq(raw_text: str, idx: int) -> dict:
    """Parse LLM output into a structured MCQ dict."""

    lines = [l.strip() for l in raw_text.strip().splitlines() if l.strip()]

    question = ""
    options = {}
    correct = ""

    for line in lines:

        if line.upper().startswith("Q:"):
            question = line[2:].strip()

        elif line.upper().startswith("A)") or line.upper().startswith("A."):
            options["A"] = line[2:].strip()

        elif line.upper().startswith("B)") or line.upper().startswith("B."):
            options["B"] = line[2:].strip()

        elif line.upper().startswith("C)") or line.upper().startswith("C."):
            options["C"] = line[2:].strip()

        elif line.upper().startswith("D)") or line.upper().startswith("D."):
            options["D"] = line[2:].strip()

        elif line.upper().startswith("CORRECT:"):
            correct = line.split(":", 1)[1].strip().upper()
            # Handle cases like "CORRECT: A)" or "CORRECT: A"
            correct = correct.replace(")", "").replace(".", "").strip()

    if not question or len(options) < 4 or correct not in ["A", "B", "C", "D"]:
        # Fallback: return a simple placeholder so exam isn't broken
        return {
            "index": idx,
            "question": question or raw_text,
            "options": options if len(options) == 4 else {"A": "Option A", "B": "Option B", "C": "Option C", "D": "Option D"},
            "correct": correct if correct in ["A", "B", "C", "D"] else "A"
        }

    return {
        "index": idx,
        "question": question,
        "options": options,
        "correct": correct
    }


# ==========================================================
# STRICT ANSWER EVALUATION
# ==========================================================

def evaluate_answer_improved(question,answer):

    prompt=f"""
You are a strict university examiner.

Rules:
- Incorrect answers → 0-3
- Partial answers → 4-6
- Mostly correct → 7-8
- Perfect answer → 9-10

Question:
{question}

Student Answer:
{answer}

Return format:

Score: X/10
Feedback: Explanation
Improvements: What student should improve
"""

    response=client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role":"user","content":prompt}],
        temperature=0.0,
        max_tokens=200
    )

    return response.choices[0].message.content


# ==========================================================
# PAGE 1 — UPLOAD
# ==========================================================

def _render_upload_page():

    st.subheader("Step 1: Upload Data Source")

    uploaded_files=st.file_uploader(
        "Upload Files",
        type=["pdf","docx","doc","pptx","ppt"],
        accept_multiple_files=True
    )

    if not uploaded_files:
        return

    total_pages=count_pages_in_files(uploaded_files)

    if total_pages==0:
        st.error("No readable content found.")
        return

    st.success(f"Total Pages Detected: {total_pages}")

    if st.button("Next ➡ Configure Exam"):

        st.session_state.uploaded_files=uploaded_files
        st.session_state.total_pages=total_pages
        st.session_state.exam_step="config"
        st.rerun()


# ==========================================================
# PAGE 2 — CONFIG
# ==========================================================

def _render_config_page():

    st.subheader("Step 2: Configure Exam")
    st.write("""
Instructions for the exam :-
1. Please note that you need to submit the answer in the specified time
2. Failure to submit the answer in specified time will lead to Zero marks for that specific question
3. please ensure time per question you set is more than enough to answer the question 
4. in case of text exam mode, you need to type the answer using your keyboard
5. in case of image exam mode , you need to write the answer on your note book with pen and upload the neat image of it to the system
6. in case of image exam mode, please maintain neat handwritting
7. in case of MCQ exam mode, you will be given 4 options for each question – select the correct one
""")
    total_pages=st.session_state.total_pages

    eval_mode=st.radio(
        "Evaluation Type",
        ["Text","Image (Handwritten)","MCQ"]
    )

    num_q=st.number_input(
        "Number of Questions",
        min_value=1,
        max_value=total_pages,
        value=min(5,total_pages)
    )

    time_limit=st.number_input(
        "Minutes per Question",
        min_value=1,
        max_value=30,
        value=2
    )
    
    if st.button("🚀 Start Examination"):

        with st.spinner("Generating questions..."):

            if eval_mode == "MCQ":
                mcq_qs = extract_mcq_questions_from_files(
                    st.session_state.uploaded_files,
                    int(num_q)
                )
                st.session_state.mcq_questions = mcq_qs
                st.session_state.questions = [m["question"] for m in mcq_qs]
            else:
                qs = extract_questions_from_files(
                    st.session_state.uploaded_files,
                    int(num_q)
                )
                st.session_state.questions = qs
                st.session_state.mcq_questions = []

        st.session_state.eval_mode=eval_mode
        st.session_state.time_limit=int(time_limit)
        st.session_state.current_q_idx=0
        st.session_state.results=[]
        st.session_state.start_time=time.time()

        st.session_state.exam_step="exam"
        st.rerun()


# ==========================================================
# EXAM PAGE
# ==========================================================

def _render_exam():

    q_idx=st.session_state.current_q_idx
    total_qs=len(st.session_state.questions)

    if q_idx>=total_qs:
        st.header("🎓 Final Exam Report")

        results=st.session_state.results
        df=pd.DataFrame(results)

        st.dataframe(df)

        if st.button("⬅ Go Back to Upload Page"):
            _reset_exam()
            st.rerun()

        return

    curr_q=st.session_state.questions[q_idx]

    st.subheader(f"Question {q_idx+1} of {total_qs}")
    st.info(curr_q)

    # TIMER
    limit_seconds=st.session_state.time_limit*60
    elapsed=int(time.time()-st.session_state.start_time)
    remaining=max(limit_seconds-elapsed,0)

    mins=remaining//60
    secs=remaining%60

    st.warning(f"⏱ Time Remaining: {mins:02d}:{secs:02d}")

    time_up=remaining<=0

    # TEXT / IMAGE / MCQ INPUT
    if st.session_state.eval_mode=="MCQ":

        mcq_data = st.session_state.mcq_questions[q_idx]
        opts = mcq_data["options"]

        radio_key = f"mcq_radio_{q_idx}"

        option_labels = [f"{key}) {val}" for key, val in opts.items()]

        selected = st.radio(
            "Select your answer:",
            option_labels,
            index=None,
            key=radio_key,
            disabled=time_up
        )

        # Extract the letter from "A) some text"
        user_ans = selected.split(")")[0].strip() if selected else ""

    elif st.session_state.eval_mode=="Image (Handwritten)":

        img=st.file_uploader(
            "Upload handwritten answer",
            type=["jpg","png","jpeg"],
            key=f"img_{q_idx}"
        )

        text_key=f"img_text_{q_idx}"

        if text_key not in st.session_state:
            st.session_state[text_key]=""

        if img and not time_up:

            with tempfile.NamedTemporaryFile(delete=False,suffix=".png") as tmp:
                tmp.write(img.read())
                path=tmp.name

            result=evaluate_image_answer(path)
            st.session_state[text_key]=result["cleaned_answer"]

            os.remove(path)

        user_ans=st.text_area(
            "Extracted Text",
            key=text_key,
            height=200,
            disabled=time_up
        )

    else:

        text_key=f"text_q_{q_idx}"

        if text_key not in st.session_state:
            st.session_state[text_key]=""

        user_ans=st.text_area(
            "Type your answer",
            key=text_key,
            height=200,
            disabled=time_up
        )

    # SUBMIT
    if st.button("Submit Answer",key=f"submit_{q_idx}"):

        if st.session_state.eval_mode == "MCQ":
            # ---- MCQ AUTO-GRADING ----
            mcq_data = st.session_state.mcq_questions[q_idx]
            correct_letter = mcq_data["correct"]
            correct_text = mcq_data["options"].get(correct_letter, "")

            if not user_ans.strip() or time_up:
                score = "Score: 0/10"
                evaluation = f"Score: 0/10\nFeedback: No answer submitted.\nCorrect Answer: {correct_letter}) {correct_text}"
            elif user_ans.strip().upper() == correct_letter:
                score = "Score: 10/10"
                evaluation = f"Score: 10/10\nFeedback: ✅ Correct! The answer is {correct_letter}) {correct_text}"
            else:
                score = "Score: 0/10"
                evaluation = f"Score: 0/10\nFeedback: ❌ Incorrect. You selected {user_ans}). The correct answer is {correct_letter}) {correct_text}"

        else:
            # ---- TEXT / IMAGE EVALUATION ----
            if not user_ans.strip():

                evaluation="Score: 0/10\nFeedback: No answer."
                score="Score: 0/10"

            else:

                with st.spinner("Evaluating..."):
                    evaluation=evaluate_answer_improved(curr_q,user_ans)

                score_line=[l for l in evaluation.split("\n") if "Score:" in l]
                score=score_line[0] if score_line else "Score: 0/10"

        st.session_state.results.append({
            "Q#":q_idx+1,
            "Question":curr_q,
            "Score":score,
            "Feedback":evaluation
        })

        st.session_state.current_q_idx+=1
        st.session_state.start_time=time.time()

        st.rerun()

    # ACTIVE TIMER LOOP
    if remaining>0:
        time.sleep(1)
        st.rerun()


# ==========================================================
# MAIN
# ==========================================================

def run_exam_mode():

    st.title("📝 Proctor Exam Mode")

    defaults = {
        "exam_step":"upload",
        "questions":[],
        "mcq_questions":[],
        "current_q_idx":0,
        "results":[],
        "start_time":None,
        "eval_mode":"Text",
        "time_limit":2
    }

    for k,v in defaults.items():
        if k not in st.session_state:
            st.session_state[k]=v

    if st.session_state.exam_step=="upload":
        _render_upload_page()

    elif st.session_state.exam_step=="config":
        _render_config_page()

    elif st.session_state.exam_step=="exam":
        _render_exam()